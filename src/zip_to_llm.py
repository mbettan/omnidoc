#!/usr/bin/env python3
"""
zip_to_llm.py — Omnidoc: Transform any ZIP archive into a single LLM-ready document.

v1.1 — Adds progress bar, token capping, PDF page cap, plain-text mode,
       and installable console entrypoint (`omnidoc`).
"""

import argparse
import csv
import io
import json
import os
import re
import sys
import zipfile
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

# ---------------------------------------------------------------------------
# Optional dependency imports (graceful degradation)
# ---------------------------------------------------------------------------
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pdf2image import convert_from_bytes
except ImportError:
    convert_from_bytes = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import docx as python_docx
except ImportError:
    python_docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak, Preformatted
    )
    from reportlab.lib.enums import TA_LEFT
except ImportError:
    SimpleDocTemplate = None

try:
    from tqdm import tqdm
except ImportError:
    tqdm = None

try:
    from google import genai
    from google.genai import types
except ImportError:
    genai = None
    types = None


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SKIP_PREFIXES = ("__MACOSX",)
SKIP_NAMES = (".DS_Store",)
MAX_UNCOMPRESSED_BYTES = 100 * 1024 * 1024  # 100 MB
DEFAULT_MAX_PDF_PAGES = 100
OCR_MIN_CONFIDENCE = 60
OCR_MIN_CHARS = 20
OCR_MIN_ALPHA_RATIO = 0.5
TEXT_PAGE_MIN_CHARS = 15
OCR_FALLBACK_CONFIDENCE = 70.0
_SAFE_PUNCTUATION = frozenset(".,;:!?()[]{}'\"-_$%=+/\\#@&*<>|~")
_HIGH_CONFIDENCE_THRESHOLD = 85
_SHORT_TEXT_MIN_CHARS = 4
CHARS_PER_TOKEN = 4  # Conservative heuristic for GPT/Claude-family tokenizers

KNOWN_TEXT_FILENAMES = {
    "readme", "license", "licence", "copying", "authors", "contributors",
    "changelog", "changes", "notice", "install", "manifest", "makefile",
    "dockerfile", "jenkinsfile", "vagrantfile", "procfile", "todo",
    "history", "news", "version",
}

FORMAT_MARKDOWN = "md"
FORMAT_TEXT = "txt"


# ---------------------------------------------------------------------------
# Token estimation
# ---------------------------------------------------------------------------
def estimate_tokens(text: str) -> int:
    """Fast heuristic: ~4 characters per token (good enough for capping)."""
    return max(1, len(text) // CHARS_PER_TOKEN)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def should_skip(path: str) -> bool:
    parts = path.split("/")
    if any(p.startswith(SKIP_PREFIXES) for p in parts):
        return True
    base = parts[-1]
    if not base:
        return True
    if base in SKIP_NAMES:
        return True
    # Ensure we don't drop traversal paths like ..
    if any(p.startswith(".") and p not in (".", "..") for p in parts if p):
        return True
    return False


def decode_bytes(data: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def looks_like_text(data: bytes, sample_size: int = 4096) -> bool:
    if not data:
        return False
    sample = data[:sample_size]
    if b"\x00" in sample:
        return False
    try:
        sample.decode("utf-8")
        return True
    except UnicodeDecodeError:
        pass
    try:
        decoded = sample.decode("latin-1")
    except UnicodeDecodeError:
        return False
    printable = sum(1 for c in decoded if c.isprintable() or c in "\r\n\t")
    return printable / max(len(decoded), 1) > 0.85


def collapse_blank_lines(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", text)


def md_table(rows: List[List[str]]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [[(str(c) if c is not None else "").replace("|", "\\|").replace("\n", " ")
             for c in (r + [""] * (width - len(r)))] for r in rows]
    header = norm[0]
    body = norm[1:] if len(norm) > 1 else []
    out = ["| " + " | ".join(header) + " |"]
    out.append("| " + " | ".join(["---"] * width) + " |")
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def plain_table(rows: List[List[str]]) -> str:
    """[FEATURE 5] Render a table as whitespace-aligned columns (no | or ---)."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [[(str(c) if c is not None else "").replace("\n", " ")
             for c in (r + [""] * (width - len(r)))] for r in rows]
    col_widths = [max(len(r[i]) for r in norm) for i in range(width)]
    lines = []
    for r in norm:
        lines.append("  ".join(c.ljust(col_widths[i]) for i, c in enumerate(r)).rstrip())
    return "\n".join(lines)


class _HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data):
        self.parts.append(data)

    def get_text(self):
        return "".join(self.parts)


def strip_html(html: str) -> str:
    s = _HTMLStripper()
    try:
        s.feed(html)
        return s.get_text()
    except Exception:
        return html


_TESSERACT_AVAILABLE: Optional[bool] = None

def is_tesseract_available() -> bool:
    global _TESSERACT_AVAILABLE
    if _TESSERACT_AVAILABLE is not None:
        return _TESSERACT_AVAILABLE
    if pytesseract is None:
        _TESSERACT_AVAILABLE = False
        return False
    try:
        pytesseract.get_tesseract_version()
        _TESSERACT_AVAILABLE = True
    except Exception:
        _TESSERACT_AVAILABLE = False
    return _TESSERACT_AVAILABLE


def _compute_alpha_ratio(text: str) -> float:
    non_space = [c for c in text if not c.isspace()]
    if not non_space:
        return 0.0
    good = sum(c.isalnum() or c in _SAFE_PUNCTUATION for c in non_space)
    return good / len(non_space)


def _is_reliable(text: str, avg_conf: float, min_confidence: int) -> bool:
    if not text:
        return False
    alpha_ratio = _compute_alpha_ratio(text)
    if avg_conf >= _HIGH_CONFIDENCE_THRESHOLD:
        return len(text) >= _SHORT_TEXT_MIN_CHARS and alpha_ratio >= 0.4
    return (
        avg_conf >= min_confidence
        and len(text) >= OCR_MIN_CHARS
        and alpha_ratio >= OCR_MIN_ALPHA_RATIO
    )


def ocr_with_confidence(image, min_confidence: int = OCR_MIN_CONFIDENCE) -> Tuple[str, float, bool]:
    if pytesseract is None or image is None:
        return "", 0.0, False

    try:
        data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        if not isinstance(data, dict):
            raise TypeError("image_to_data did not return a dictionary")
        confidences = [int(c) for c in data.get("conf", []) if c != -1 and c != "-1"]
        words = [w for w in data.get("text", []) if w and str(w).strip()]
        text = " ".join(words).strip()

        if not confidences or not text:
            return "", 0.0, False

        avg_conf = sum(confidences) / len(confidences)
        reliable = _is_reliable(text, avg_conf, min_confidence)
        return text, avg_conf, reliable
    except (TypeError, KeyError, ValueError, IndexError, EnvironmentError):
        text = pytesseract.image_to_string(image).strip()
        if text:
            alpha_ratio = _compute_alpha_ratio(text)
            if len(text) >= OCR_MIN_CHARS and alpha_ratio >= OCR_MIN_ALPHA_RATIO:
                return text, OCR_FALLBACK_CONFIDENCE, True
        return "", 0.0, False


def _has_significant_vectors(page) -> bool:
    rect_count = len(getattr(page, "rects", []) or [])
    line_count = len(getattr(page, "lines", []) or [])
    curve_count = len(getattr(page, "curves", []) or [])
    # Weight curves heavier: a single chart often has 50+ curves but few rects
    weighted = rect_count + line_count + (curve_count * 3)
    return weighted > 15


def classify_pdf_pages(pdf_data: bytes, min_chars: int = TEXT_PAGE_MIN_CHARS) -> List[Dict]:
    results = []
    if pdfplumber is None:
        return results
    try:
        with pdfplumber.open(io.BytesIO(pdf_data)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                try:
                    text = ""
                    try:
                        text = (page.extract_text() or "").strip()
                    except Exception:
                        pass
                    has_text = len(text) >= min_chars
                    
                    has_visuals = bool(getattr(page, "images", None)) or _has_significant_vectors(page)

                    if has_text and not has_visuals:
                        kind = "text"
                    elif has_text and has_visuals:
                        kind = "mixed"
                    else:
                        kind = "image_only"
                    results.append({"page_num": i, "kind": kind, "text": text})
                except Exception:
                    results.append({"page_num": i, "kind": "image_only", "text": ""})
    except Exception as e:
        print(f"Warning: pdfplumber failed to open PDF: {e}", file=sys.stderr)
    return results


def _contiguous_runs(nums):
    if not nums:
        return []
    runs, start, prev = [], nums[0], nums[0]
    for n in nums[1:]:
        if n == prev + 1:
            prev = n
        else:
            runs.append((start, prev))
            start = prev = n
    runs.append((start, prev))
    return runs


class VisualBundle:
    """Accumulates pages/images that couldn't be OCR'd into a single PDF."""

    def __init__(self, max_dim: int = 1600, quality: int = 75):
        self._entries = []  # list of (anchor, source, page_num, data, kind)
        self.max_dim = max_dim
        self.quality = quality

    def add_pdf_page(self, source: str, page_num: int, page_bytes: bytes) -> Tuple[str, int]:
        from pypdf import PdfReader, PdfWriter
        bundle_page = len(self._entries) + 1
        anchor = f"visual-{bundle_page:04d}"
        try:
            reader = PdfReader(io.BytesIO(page_bytes))
            if page_num - 1 >= len(reader.pages):
                print(f"Warning: anchor {anchor} references page {page_num} of {source}, "
                      f"but source has only {len(reader.pages)} pages — entry dropped", file=sys.stderr)
                return anchor, bundle_page
            
            writer = PdfWriter()
            writer.add_page(reader.pages[page_num - 1])
            buf = io.BytesIO()
            writer.write(buf)
            single_page_bytes = buf.getvalue()
        except Exception as e:
            print(f"Warning: failed to extract page {page_num} from {source}: {e} — entry dropped", file=sys.stderr)
            return anchor, bundle_page
        self._entries.append((anchor, source, page_num, single_page_bytes, "pdf"))
        return anchor, bundle_page

    def add_image(self, source: str, image_data: bytes) -> Tuple[str, int]:
        bundle_page = len(self._entries) + 1
        anchor = f"visual-{bundle_page:04d}"
        optimized_data = self._optimize(image_data)
        self._entries.append((anchor, source, None, optimized_data, "image"))
        return anchor, bundle_page

    def _optimize(self, image_data: bytes) -> bytes:
        if Image is None:
            return image_data
        try:
            img = Image.open(io.BytesIO(image_data))
            if img.mode not in ("RGB", "L", "CMYK"):
                img = img.convert("RGB")
            w, h = img.size
            scale = min(1.0, self.max_dim / max(w, h))
            if scale < 1.0:
                try:
                    resample = Image.Resampling.LANCZOS
                except AttributeError:
                    resample = Image.LANCZOS
                img = img.resize((int(w * scale), int(h * scale)), resample)
            
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=self.quality, optimize=True)
            return buf.getvalue()
        except Exception:
            return image_data

    def write(self, output_path: Path) -> int:
        if not self._entries:
            return 0
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError:
            raise RuntimeError("pypdf is required to write the visual bundle.")

        writer = PdfWriter()

        for anchor, source, page_num, data, kind in self._entries:
            try:
                if kind == "pdf":
                    reader = PdfReader(io.BytesIO(data))
                    writer.add_page(reader.pages[0])
                else:  # image
                    img = Image.open(io.BytesIO(data))
                    if img.mode != "RGB":
                        img = img.convert("RGB")
                    pdf_buf = io.BytesIO()
                    img.save(pdf_buf, format="PDF")
                    pdf_buf.seek(0)
                    
                    img_reader = PdfReader(pdf_buf)
                    writer.add_page(img_reader.pages[0])
            except Exception as e:
                print(f"Warning: failed to write visual bundle entry '{anchor}': {e}", file=sys.stderr)

        if len(writer.pages) > 0:
            try:
                for page in writer.pages:
                    try:
                        page.compress_content_streams()
                    except Exception:
                        pass
                with open(output_path, "wb") as f:
                    writer.write(f)
                return len(self._entries)
            except Exception as e:
                print(f"Error: failed to write visual bundle file {output_path}: {e}", file=sys.stderr)
        return 0

    def __len__(self):
        return len(self._entries)


# ---------------------------------------------------------------------------
# Output formatter — abstracts Markdown vs plain text
# ---------------------------------------------------------------------------
class OutputFormatter:
    """[FEATURE 5] Switches markdown decorations on/off based on output format."""

    def __init__(self, fmt: str = FORMAT_MARKDOWN):
        self.fmt = fmt
        self.is_md = (fmt == FORMAT_MARKDOWN)

    def title(self, text: str) -> str:
        return f"# {text}" if self.is_md else text.upper()

    def h2(self, text: str) -> str:
        return f"## {text}" if self.is_md else f"\n=== {text} ==="

    def h3(self, text: str) -> str:
        return f"### {text}" if self.is_md else f"\n--- {text} ---"

    def h4(self, text: str) -> str:
        return f"#### {text}" if self.is_md else f"\n[{text}]"

    def bold_kv(self, key: str, value: str) -> str:
        return f"**{key}:** {value}" if self.is_md else f"{key}: {value}"

    def code(self, text: str, lang: str = "") -> str:
        if self.is_md:
            return f"```{lang}\n{text}\n```"
        return text

    def inline_code(self, text: str) -> str:
        return f"`{text}`" if self.is_md else text

    def separator(self) -> str:
        return "---" if self.is_md else "\n" + ("=" * 60) + "\n"

    def table(self, rows: List[List[str]]) -> str:
        return md_table(rows) if self.is_md else plain_table(rows)

    def bullet(self, text: str) -> str:
        return f"- {text}" if self.is_md else f"  • {text}"

    def italic_note(self, text: str) -> str:
        return f"_[{text}]_" if self.is_md else f"[{text}]"


# ---------------------------------------------------------------------------
# Extractors — every extractor receives (data, name, fmt, opts) and returns str
# ---------------------------------------------------------------------------
class ExtractOptions:
    """Bag of per-run extraction parameters."""
    def __init__(self, max_pdf_pages: int = DEFAULT_MAX_PDF_PAGES,
                 isolate_non_text: bool = False,
                 ocr_min_confidence: int = OCR_MIN_CONFIDENCE,
                 visual_max_dim: int = 1600,
                 vision_fallback: str = "none",
                 tesseract_available: bool = True,
                 text_page_min_chars: int = TEXT_PAGE_MIN_CHARS,
                 ocr_render_dpi: int = 200,
                 visual_quality: int = 75):
        self.max_pdf_pages = max_pdf_pages
        self.isolate_non_text = isolate_non_text
        self.ocr_min_confidence = ocr_min_confidence
        self.visual_max_dim = visual_max_dim
        self.vision_fallback = vision_fallback
        self.tesseract_available = tesseract_available
        self.text_page_min_chars = text_page_min_chars
        self.ocr_render_dpi = ocr_render_dpi
        self.visual_quality = visual_quality
        self.visual_bundle: Optional[VisualBundle] = None


def extract_pdf(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")
    if pdfplumber is None:
        return fmt.italic_note("pdfplumber not installed; PDF skipped.")

    out: List[str] = []
    truncated = False
    total_pages = 0

    try:
        pages_info = classify_pdf_pages(data, min_chars=opts.text_page_min_chars)
        if not pages_info:
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(data))
                if len(reader.pages) > 0:
                    return fmt.italic_note(
                        f"PDF has {len(reader.pages)} pages but could not be classified "
                        "(pdfplumber parse error)"
                    )
            except Exception:
                pass
            return fmt.italic_note("Empty or unparseable PDF")
        total_pages = len(pages_info)
        limit = min(total_pages, opts.max_pdf_pages)
        if total_pages > limit:
            truncated = True

        needs_render_pnos = [p["page_num"] for p in pages_info[:limit] if p["kind"] != "text"]
        rendered_images = {}

        if needs_render_pnos and convert_from_bytes and pytesseract:
            for start, end in _contiguous_runs(sorted(needs_render_pnos)):
                try:
                    imgs = convert_from_bytes(data, first_page=start, last_page=end, dpi=opts.ocr_render_dpi)
                    for offset, img in enumerate(imgs):
                        rendered_images[start + offset] = img
                except Exception as e:
                    print(f"Warning: failed to render pages {start}-{end} of {name}: {e}", file=sys.stderr)

        for p in pages_info[:limit]:
            pno = p["page_num"]
            section = []

            try:
                section.append(fmt.h3(f"Page {pno}"))
                if p["kind"] == "text":
                    if p["text"]:
                        section.append(p["text"])
                    else:
                        section.append(fmt.italic_note("empty page"))
                    out.append("\n\n".join(section))
                    continue

                img = rendered_images.get(pno)
                ocr_text = ""
                avg_conf = 0.0
                reliable = False

                if img is not None:
                    try:
                        ocr_text, avg_conf, reliable = ocr_with_confidence(img, min_confidence=opts.ocr_min_confidence)
                    except Exception:
                        pass

                if reliable:
                    if p["kind"] == "mixed" and p["text"]:
                        section.append(p["text"])
                    section.append(f"\n*[{fmt.bold_kv('OCR text', f'confidence {avg_conf:.0f}%')}]*\n{ocr_text}")
                else:
                    if p["kind"] == "mixed" and p["text"]:
                        section.append(p["text"])
                    if opts.isolate_non_text and opts.visual_bundle is not None:
                        anchor, bundle_page = opts.visual_bundle.add_pdf_page(name, pno, data)
                        section.append(
                            f"> ⚠️ **Visual content** — see anchor `{fmt.inline_code(anchor)}` "
                            f"(companion PDF page {bundle_page}, source: `{fmt.inline_code(name)}`, source page {pno})"
                        )
                    else:
                        section.append(fmt.italic_note("no extractable text"))
            except Exception as pe:
                section.append(fmt.italic_note(f"Page extraction error: {pe}"))

            out.append("\n\n".join(section))

    except Exception as e:
        return fmt.italic_note(f"PDF parse error: {e}")

    if truncated:
        out.append(fmt.italic_note(
            f"PDF truncated: showing {opts.max_pdf_pages} of {total_pages} pages"))

    return "\n\n".join(out) if out else fmt.italic_note("Empty PDF.")


def _bundle_image(name: str, data: bytes, fmt: OutputFormatter, opts: ExtractOptions, reason: str = "") -> str:
    anchor, bundle_page = opts.visual_bundle.add_image(name, data)
    suffix = f" ({reason})" if reason else ""
    return (
        f"> ⚠️ **Image with no extractable text**{suffix} — see anchor `{fmt.inline_code(anchor)}` "
        f"(companion PDF page {bundle_page}, source: `{fmt.inline_code(name)}`)"
    )


def extract_image(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                  opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")

    if Image is None or pytesseract is None or not opts.tesseract_available:
        if opts.isolate_non_text and opts.visual_bundle is not None:
            return _bundle_image(name, data, fmt, opts, "Pillow/pytesseract missing or unconfigured")
        if Image is None or pytesseract is None:
            return fmt.italic_note("Pillow/pytesseract not installed; image OCR skipped.")

    text = ""
    conf = 0.0
    reliable = False

    if Image is not None and pytesseract is not None and opts.tesseract_available:
        try:
            img = Image.open(io.BytesIO(data))
        except Exception as e:
            if opts.isolate_non_text and opts.visual_bundle is not None:
                return _bundle_image(name, data, fmt, opts, f"Image open error: {e}")
            return fmt.italic_note(f"Image OCR error: {e}")

        try:
            text, conf, reliable = ocr_with_confidence(img, min_confidence=opts.ocr_min_confidence)
        except Exception as e:
            if opts.isolate_non_text and opts.visual_bundle is not None:
                return _bundle_image(name, data, fmt, opts, f"Image OCR error: {e}")
            return fmt.italic_note(f"Image OCR error: {e}")

    if reliable:
        return f"*[OCR, confidence {conf:.0f}%]*\n\n{text}"

    if opts.isolate_non_text and opts.visual_bundle is not None:
        return _bundle_image(name, data, fmt, opts)

    if opts.vision_fallback == "gemini" and genai is not None and types is not None:
        project_id = os.getenv("PROJECT_ID")
        if not project_id:
            return fmt.italic_note("no text detected by OCR (vision fallback disabled: set PROJECT_ID)")
        try:
            import mimetypes
            img_type, _ = mimetypes.guess_type(name)
            if not img_type:
                img_type = "image/jpeg"
            client = genai.Client(vertexai=True, project=project_id, location="global")
            image_part = types.Part.from_bytes(data=data, mime_type=img_type)
            response = client.models.generate_content(
                model=os.getenv("OMNIDOC_VISION_MODEL", "gemini-3.1-pro-preview"),
                contents=[image_part, "Analyze this image in detail and describe its visual contents."]
            )
            gemini_text = (response.text or "").strip()
            if gemini_text:
                return f"{fmt.italic_note('Image Description (Gemini)')}\n\n{gemini_text}"
        except Exception as ge:
            return fmt.italic_note(f"no text detected by OCR (vision fallback error: {ge})")

    return fmt.italic_note("no text detected by OCR")


def extract_excel(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                  opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")
    if openpyxl is None:
        return fmt.italic_note("openpyxl not installed; Excel skipped.")
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception as e:
        return fmt.italic_note(f"Excel parse error: {e}")

    sections = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(["" if v is None else str(v) for v in row])
        while rows and not any(c.strip() for c in rows[-1]):
            rows.pop()
        body = fmt.table(rows) if rows else fmt.italic_note("empty sheet")
        sections.append(f"{fmt.h4(f'Sheet: {sheet_name}')}\n\n{body}")
    return "\n\n".join(sections) if sections else fmt.italic_note("no sheets")


def extract_csv(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                opts: Optional[ExtractOptions] = None, delimiter: str = ",") -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")
    text = decode_bytes(data)
    try:
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        rows = [r for r in reader]
        if not rows:
            return fmt.italic_note("empty file")
        return fmt.table(rows)
    except Exception as e:
        return fmt.italic_note(f"CSV parse error: {e}")


def extract_tsv(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()
    return extract_csv(data, name, fmt, opts, delimiter="\t")


def extract_docx(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                 opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")
    if python_docx is None:
        return fmt.italic_note("python-docx not installed; DOCX skipped.")
    try:
        doc = python_docx.Document(io.BytesIO(data))
    except Exception as e:
        return fmt.italic_note(f"DOCX parse error: {e}")

    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            style = (para.style.name or "").lower() if para.style else ""
            if "heading 1" in style:
                parts.append(fmt.h3(t))
            elif "heading" in style:
                parts.append(fmt.h4(t))
            else:
                parts.append(t)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        if rows:
            parts.append(fmt.table(rows))
    return "\n\n".join(parts) if parts else fmt.italic_note("empty document")


def extract_pptx(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                 opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")
    if Presentation is None:
        return fmt.italic_note("python-pptx not installed; PPTX skipped.")
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        return fmt.italic_note(f"PPTX parse error: {e}")

    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks = [fmt.h4(f"Slide {i}")]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text.strip())
        sections.append("\n\n".join(c for c in chunks if c))
    return "\n\n".join(sections) if sections else fmt.italic_note("empty presentation")


def extract_text(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                 opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()

    if not data:
        return fmt.italic_note("empty file")
    text = decode_bytes(data)
    ext = Path(name).suffix.lower()
    if ext == ".json":
        try:
            obj = json.loads(text)
            text = json.dumps(obj, indent=2, ensure_ascii=False)
        except Exception:
            pass
        return fmt.code(text, "json")
    if ext in (".html", ".htm"):
        return strip_html(text)
    if ext == ".md":
        return text
    if ext in (".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".rb",
               ".go", ".rs", ".sh", ".sql", ".yaml", ".yml", ".toml", ".ini"):
        lang = ext.lstrip(".")
        lang_map = {"py": "python", "js": "javascript", "ts": "typescript", "rb": "ruby", "sh": "bash"}
        lang = lang_map.get(lang, lang)
        return fmt.code(text, lang)
    if ext == ".xml":
        return fmt.code(text, "xml")
    return text


def extract_unsupported(data: bytes, name: str, fmt: Optional[OutputFormatter] = None,
                        opts: Optional[ExtractOptions] = None) -> str:
    fmt = fmt or OutputFormatter()
    opts = opts or ExtractOptions()
    ext = Path(name).suffix.lower()
    ext_label = ext if ext else "no extension"
    return fmt.italic_note(
        f"unsupported file type ({ext_label}) — skipped")


# ---------------------------------------------------------------------------
# Format router
# ---------------------------------------------------------------------------
EXT_MAP: Dict[str, Tuple[str, Callable]] = {
    ".pdf":  ("PDF Document", extract_pdf),
    ".jpg":  ("Image", extract_image),
    ".jpeg": ("Image", extract_image),
    ".png":  ("Image", extract_image),
    ".bmp":  ("Image", extract_image),
    ".tiff": ("Image", extract_image),
    ".tif":  ("Image", extract_image),
    ".gif":  ("Image", extract_image),
    ".webp": ("Image", extract_image),
    ".xlsx": ("Excel Spreadsheet", extract_excel),
    ".xlsm": ("Excel Spreadsheet", extract_excel),
    ".xls":  ("Excel Spreadsheet", extract_excel),
    ".csv":  ("CSV", extract_csv),
    ".tsv":  ("TSV", extract_tsv),
    ".docx": ("Word Document", extract_docx),
    ".doc":  ("Word Document", extract_docx),
    ".pptx": ("PowerPoint", extract_pptx),
    ".ppt":  ("PowerPoint", extract_pptx),
    ".txt":  ("Text", extract_text),
    ".md":   ("Markdown", extract_text),
    ".json": ("JSON", extract_text),
    ".xml":  ("XML", extract_text),
    ".html": ("HTML", extract_text),
    ".htm":  ("HTML", extract_text),
    ".py":   ("Python Source", extract_text),
    ".js":   ("JavaScript Source", extract_text),
    ".ts":   ("TypeScript Source", extract_text),
    ".java": ("Java Source", extract_text),
    ".c":    ("C Source", extract_text),
    ".cpp":  ("C++ Source", extract_text),
    ".h":    ("C Header", extract_text),
    ".rb":   ("Ruby Source", extract_text),
    ".go":   ("Go Source", extract_text),
    ".rs":   ("Rust Source", extract_text),
    ".sh":   ("Shell Script", extract_text),
    ".sql":  ("SQL", extract_text),
    ".yaml": ("YAML", extract_text),
    ".yml":  ("YAML", extract_text),
    ".toml": ("TOML", extract_text),
    ".ini":  ("INI", extract_text),
    # Media files mapped explicitly for graceful skipping
    ".mp4":  ("Video", extract_unsupported),
    ".m4a":  ("Audio", extract_unsupported),
    ".mp3":  ("Audio", extract_unsupported),
    ".wav":  ("Audio", extract_unsupported),
    ".mov":  ("Video", extract_unsupported),
    ".avi":  ("Video", extract_unsupported),
    ".mkv":  ("Video", extract_unsupported),
    ".webm": ("Video", extract_unsupported),
}


def route(name: str,
          peek: Optional[Callable[[int], bytes]] = None
          ) -> Tuple[str, Callable]:
    p = Path(name)
    ext = p.suffix.lower()
    if ext in EXT_MAP:
        return EXT_MAP[ext]
    # Check if base name or stem matches known text filenames (e.g. README, Dockerfile.dev)
    stem_lower = p.stem.lower() if p.stem else p.name.lower()
    head = stem_lower.split(".")[0]
    if head in KNOWN_TEXT_FILENAMES or p.name.lower().split(".")[0] in KNOWN_TEXT_FILENAMES:
        return ("Text", extract_text)
    if peek is not None:
        try:
            sample = peek(4096)
            if looks_like_text(sample):
                return ("Text", extract_text)
        except Exception:
            pass
    return ("Unsupported", extract_unsupported)


# ---------------------------------------------------------------------------
# Progress bar abstraction
# ---------------------------------------------------------------------------
class ProgressReporter:
    """[FEATURE 2] Wraps tqdm with a graceful fallback when unavailable."""

    def __init__(self, total: int, enabled: bool = True, desc: str = "Processing"):
        self.total = total
        self.enabled = enabled and tqdm is not None and sys.stderr.isatty()
        self.bar = None
        if self.enabled:
            self.bar = tqdm(total=total, desc=desc, unit="file",
                            file=sys.stderr, leave=False,
                            bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} "
                                       "[{elapsed}<{remaining}] {postfix}")
        elif enabled:
            # Fallback line printer
            print(f"Processing {total} files…", file=sys.stderr)

    def update(self, filename: str):
        if self.bar:
            self.bar.set_postfix_str(filename[:40], refresh=False)
            self.bar.update(1)

    def close(self):
        if self.bar:
            self.bar.close()


# ---------------------------------------------------------------------------
# ZIP walker / assembler
# ---------------------------------------------------------------------------
def _format_size(n: int) -> str:
    return f"{n:,} bytes"


def process_zip(zip_path: Path,
                output_format: str = FORMAT_MARKDOWN,
                max_uncompressed: int = MAX_UNCOMPRESSED_BYTES,
                max_pdf_pages: int = DEFAULT_MAX_PDF_PAGES,
                max_tokens: Optional[int] = None,
                show_progress: bool = True,
                isolate_non_text: bool = False,
                non_text_pdf_path: Optional[Path] = None,
                ocr_min_confidence: int = OCR_MIN_CONFIDENCE,
                visual_max_dim: int = 1600,
                vision_fallback: str = "none",
                text_page_min_chars: int = TEXT_PAGE_MIN_CHARS,
                ocr_render_dpi: int = 200,
                visual_quality: int = 75) -> str:
    if not zipfile.is_zipfile(zip_path):
        raise ValueError(f"Not a valid ZIP archive: {zip_path}")

    tesseract_ok = is_tesseract_available()

    if pytesseract is not None and not tesseract_ok:
        print("WARNING: pytesseract is installed but the Tesseract binary is not "
              "available on PATH. Image OCR will be skipped; all visual content "
              "will be bundled without text extraction. Install via:\n"
              "  macOS:  brew install tesseract\n"
              "  Linux:  apt install tesseract-ocr", file=sys.stderr)

    if isolate_non_text:
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError:
            raise ValueError("pypdf must be installed to use --isolate.")
        if Image is None:
            raise ValueError("Pillow must be installed to use --isolate.")

    fmt = OutputFormatter(output_format)
    opts = ExtractOptions(
        max_pdf_pages=max_pdf_pages,
        isolate_non_text=isolate_non_text,
        ocr_min_confidence=ocr_min_confidence,
        visual_max_dim=visual_max_dim,
        vision_fallback=vision_fallback,
        tesseract_available=tesseract_ok,
        text_page_min_chars=text_page_min_chars,
        ocr_render_dpi=ocr_render_dpi,
        visual_quality=visual_quality
    )

    if isolate_non_text:
        opts.visual_bundle = VisualBundle(max_dim=visual_max_dim, quality=visual_quality)

    # Pre-process ZIP to collect recursive directory and file manifest metadata
    manifest_rows = [["Entry Path", "Type", "Compressed Size", "Uncompressed Size", "Date Modified"]]
    
    with zipfile.ZipFile(zip_path, "r") as zf:
        # Walk the entire zip recursively including folders and ignored/dotfiles
        all_infolist = zf.infolist()
        # Sort alphabetically by filename
        all_infolist = sorted(all_infolist, key=lambda x: x.filename)
        
        for info in all_infolist:
            entry_path = info.filename
            
            # Determine entry type
            if entry_path.endswith("/"):
                entry_type = "Folder"
            else:
                # Determine if it's a file or a special skipped prefix/name
                is_skipped = should_skip(entry_path)
                if is_skipped:
                    entry_type = "File (Skipped Metadata)"
                else:
                    # Peak the format route
                    def _peek_temp(n_bytes: int, _name=entry_path, _zf=zf) -> bytes:
                        with _zf.open(_name, "r") as fh:
                            return fh.read(n_bytes)
                    lbl, _ = route(entry_path, peek=_peek_temp)
                    entry_type = f"File ({lbl})"
            
            # Format timestamp
            # info.date_time is a tuple: (year, month, day, hour, minute, second)
            dt_str = f"{info.date_time[0]:04d}-{info.date_time[1]:02d}-{info.date_time[2]:02d} {info.date_time[3]:02d}:{info.date_time[4]:02d}:{info.date_time[5]:02d}"
            
            manifest_rows.append([
                entry_path,
                entry_type,
                _format_size(info.compress_size),
                _format_size(info.file_size),
                dt_str
            ])

    sections: List[str] = []
    file_entries: List[Tuple[str, int]] = []
    truncated_at: Optional[str] = None
    running_tokens = 0

    with zipfile.ZipFile(zip_path, "r") as zf:
        names = [n for n in zf.namelist() if not should_skip(n)]
        names = [n for n in names if not n.endswith("/")]
        names.sort()

        progress = ProgressReporter(total=len(names), enabled=show_progress)

        try:
            for name in names:
                progress.update(name)

                try:
                    info = zf.getinfo(name)
                except KeyError:
                    sections.append(_render_error_section(name, "Unknown", "?",
                                                          "missing zip entry", fmt))
                    continue

                declared_size = info.file_size

                if declared_size > max_uncompressed:
                    file_entries.append((name, declared_size))
                    sections.append(_render_section(
                        name, "Skipped", _format_size(declared_size),
                        fmt.italic_note(
                            f"file exceeds {_format_size(max_uncompressed)} cap — skipped"),
                        fmt))
                    continue

                def _peek(n_bytes: int, _name=name, _zf=zf) -> bytes:
                    with _zf.open(_name, "r") as fh:
                        return fh.read(n_bytes)

                type_label, extractor = route(name, peek=_peek)

                if extractor is extract_unsupported:
                    file_entries.append((name, declared_size))
                    sections.append(_render_section(
                        name, type_label, _format_size(declared_size),
                        extract_unsupported(b"", name, fmt, opts), fmt))
                    continue

                try:
                    with zf.open(name, "r") as fh:
                        data = fh.read(max_uncompressed + 1)
                    if len(data) > max_uncompressed:
                        file_entries.append((name, len(data)))
                        sections.append(_render_section(
                            name, "Skipped",
                            f"> {_format_size(max_uncompressed)}",
                            fmt.italic_note(
                                "actual size exceeds cap during read — skipped"),
                            fmt))
                        continue
                except Exception as e:
                    sections.append(_render_error_section(
                        name, type_label, "?", f"Read error: {e}", fmt))
                    continue

                size = len(data)
                file_entries.append((name, size))

                try:
                    content = extractor(data, name, fmt, opts)
                except Exception as e:
                    content = fmt.italic_note(f"Extraction error: {e}")

                section = _render_section(name, type_label,
                                          _format_size(size), content, fmt)

                if max_tokens is not None:
                    section_tokens = estimate_tokens(section)
                    if running_tokens + section_tokens > max_tokens:
                        truncated_at = name
                        break
                    running_tokens += section_tokens

                sections.append(section)
        finally:
            progress.close()

    # Header / TOC
    header_lines = [fmt.title("LLM Document Bundle"), ""]
    header_lines.append(fmt.bold_kv("Source archive", fmt.inline_code(zip_path.name)))
    header_lines.append(fmt.bold_kv("Files found", str(len(file_entries))))
    if max_tokens is not None:
        header_lines.append(fmt.bold_kv("Token budget", f"{max_tokens:,}"))
        header_lines.append(fmt.bold_kv("Estimated tokens used", f"{running_tokens:,}"))
    header_lines.append("")
    
    # Beautiful documented Archive Directory Manifest table
    header_lines.append(fmt.h2("Archive Directory Manifest"))
    header_lines.append("Below is a complete recursive manifest listing all folders, files, and metadata found within the source ZIP archive (including hidden files and skipped metadata elements):")
    header_lines.append("")
    header_lines.append(fmt.table(manifest_rows))
    header_lines.append("")
    
    header_lines.append(fmt.h2("Table of Contents"))
    header_lines.append("")
    for name, _ in file_entries:
        header_lines.append(fmt.bullet(fmt.inline_code(name)))
    header_lines.append("")

    if truncated_at:
        header_lines.append(fmt.italic_note(
            f"⚠ Output truncated at file '{truncated_at}' to respect "
            f"--max-tokens {max_tokens:,}."))
        header_lines.append("")

    if isolate_non_text and non_text_pdf_path and opts.visual_bundle and len(opts.visual_bundle) > 0:
        opts.visual_bundle.write(non_text_pdf_path)

    full = "\n".join(header_lines) + "\n" + fmt.separator() + "\n\n" + \
           ("\n\n" + fmt.separator() + "\n\n").join(sections) + "\n"
    return collapse_blank_lines(full)


def _render_section(name: str, type_label: str, size_str: str,
                    content: str, fmt: OutputFormatter) -> str:
    sep = " | " if fmt.is_md else "  "
    return (
        f"{fmt.h2(f'File: {fmt.inline_code(name)}')}\n\n"
        f"{fmt.bold_kv('Type', type_label)}{sep}{fmt.bold_kv('Size', size_str)}\n\n"
        f"{content}"
    )


def _render_error_section(name: str, type_label: str, size_str: str,
                          error: str, fmt: OutputFormatter) -> str:
    return _render_section(name, type_label, size_str,
                           fmt.italic_note(error), fmt)


# ---------------------------------------------------------------------------
# PDF rendering
# ---------------------------------------------------------------------------
def _xml_escape(t: str) -> str:
    return (t.replace("&", "&amp;")
             .replace("<", "&lt;")
             .replace(">", "&gt;"))


def render_pdf(markdown_text: str, pdf_path: Path) -> None:
    if SimpleDocTemplate is None:
        raise RuntimeError("reportlab not installed; cannot render PDF.")

    doc = SimpleDocTemplate(
        str(pdf_path), pagesize=LETTER,
        leftMargin=0.75 * inch, rightMargin=0.75 * inch,
        topMargin=0.75 * inch, bottomMargin=0.75 * inch,
    )
    styles = getSampleStyleSheet()
    body = ParagraphStyle("Body", parent=styles["BodyText"],
                          fontName="Helvetica", fontSize=10, leading=13,
                          alignment=TA_LEFT)
    h1, h2, h3 = styles["Heading1"], styles["Heading2"], styles["Heading3"]
    code_style = ParagraphStyle("Code", parent=styles["Code"],
                                fontName="Courier", fontSize=8, leading=10)

    flowables = []
    in_code = False
    code_buffer: List[str] = []

    def flush_code():
        nonlocal code_buffer
        if code_buffer:
            text = "\n".join(_xml_escape(line) for line in code_buffer)
            try:
                flowables.append(Preformatted(text, code_style))
            except Exception as e:
                flowables.append(Paragraph(
                    f"<i>[code block render error: {_xml_escape(str(e))}]</i>", body))
            flowables.append(Spacer(1, 6))
            code_buffer = []

    for line in markdown_text.splitlines():
        if line.startswith("```"):
            if in_code:
                flush_code()
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buffer.append(line)
            continue
        stripped = line.rstrip()
        if not stripped:
            flowables.append(Spacer(1, 6))
            continue
        try:
            if stripped.startswith("# "):
                flowables.append(Paragraph(_xml_escape(stripped[2:]), h1))
            elif stripped.startswith("## "):
                flowables.append(Paragraph(_xml_escape(stripped[3:]), h2))
            elif stripped.startswith("### "):
                flowables.append(Paragraph(_xml_escape(stripped[4:]), h3))
            elif stripped.startswith("#### "):
                flowables.append(Paragraph(_xml_escape(stripped[5:]), h3))
            elif stripped == "---":
                flowables.append(Spacer(1, 12))
                flowables.append(PageBreak())
            else:
                flowables.append(Paragraph(_xml_escape(stripped), body))
        except Exception:
            continue

    flush_code()
    doc.build(flowables)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(
        prog="omnidoc",
        description=("Omnidoc — Transform any ZIP archive into a single "
                     "LLM-ready document."),
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    parser.add_argument("zip_file", help="Path to input .zip archive")
    parser.add_argument("-o", "--output", default=None,
                        help="Base name for output file(s) (no extension)")
    parser.add_argument("--pdf", action="store_true",
                        help="Also render and save a PDF version")
    parser.add_argument("--format", choices=[FORMAT_MARKDOWN, FORMAT_TEXT],
                        default=FORMAT_MARKDOWN,
                        help="Output format: 'md' (rich) or 'txt' (token-minimal)")
    parser.add_argument("--max-file-mb", type=int, default=100,
                        help="Per-file uncompressed size cap in MB")
    parser.add_argument("--max-pdf-pages", type=int, default=DEFAULT_MAX_PDF_PAGES,
                        help="Maximum number of pages extracted per PDF")
    parser.add_argument("--max-tokens", type=int, default=None,
                        help="Cap output at approximately N tokens "
                             "(~4 chars/token). Truncates with notice if exceeded.")
    parser.add_argument("--no-progress", action="store_true",
                        help="Disable progress bar")
    parser.add_argument("--no-visual-bundle", action="store_true",
                        help="Disable the companion visual bundle PDF")
    parser.add_argument("--ocr-min-confidence", type=int, default=OCR_MIN_CONFIDENCE,
                        help="Tesseract confidence threshold (0-100) below which content is bundled")
    parser.add_argument("--visual-max-dim", type=int, default=1600,
                        help="Max pixel dimension for embedded visual images")
    parser.add_argument("--visual-quality", type=int, default=75,
                        help="JPEG compression quality (10-100) for visual bundle images")
    parser.add_argument("--vision-fallback", choices=["none", "gemini"], default="none",
                        help="Optional cloud vision fallback model when OCR fails")
    parser.add_argument("--text-page-min-chars", type=int, default=TEXT_PAGE_MIN_CHARS,
                        help="Minimum characters for a PDF page to be classified as text-only")
    parser.add_argument("--ocr-render-dpi", type=int, default=200,
                        help="DPI resolution for selective PDF page rendering")
    parser.add_argument("-q", "--quiet", action="store_true",
                        help="Suppress non-error output")
    parser.add_argument("-V", "--version", action="version",
                        version="omnidoc 1.1.0")

    args = parser.parse_args(argv)

    zip_path = Path(args.zip_file)
    if not zip_path.exists():
        print(f"error: input file not found: {zip_path}", file=sys.stderr)
        return 1
    if not zip_path.is_file():
        print(f"error: not a file: {zip_path}", file=sys.stderr)
        return 1

    output_base = args.output or f"{zip_path.stem}_llm"
    ext = ".md" if args.format == FORMAT_MARKDOWN else ".txt"
    out_path = Path(output_base + ext)
    max_bytes = max(1, args.max_file_mb) * 1024 * 1024
    isolate = not args.no_visual_bundle
    non_text_pdf_path = Path(output_base + "_visual.pdf") if isolate else None

    try:
        content = process_zip(
            zip_path,
            output_format=args.format,
            max_uncompressed=max_bytes,
            max_pdf_pages=max(1, args.max_pdf_pages),
            max_tokens=args.max_tokens,
            show_progress=(not args.no_progress and not args.quiet),
            isolate_non_text=isolate,
            non_text_pdf_path=non_text_pdf_path,
            ocr_min_confidence=args.ocr_min_confidence,
            visual_max_dim=args.visual_max_dim,
            vision_fallback=args.vision_fallback,
            text_page_min_chars=args.text_page_min_chars,
            ocr_render_dpi=args.ocr_render_dpi,
            visual_quality=args.visual_quality,
        )
    except ValueError as e:
        print(f"error: {e}", file=sys.stderr)
        return 1
    except Exception as e:
        print(f"error: unexpected failure: {e}", file=sys.stderr)
        return 1

    try:
        out_path.write_text(content, encoding="utf-8")
        if not args.quiet:
            tokens = estimate_tokens(content)
            print(f"✓ {args.format.upper()} written: {out_path} "
                  f"(~{tokens:,} tokens)", file=sys.stderr)
            if isolate and non_text_pdf_path and non_text_pdf_path.exists():
                size_bytes = non_text_pdf_path.stat().st_size
                print(f"✓ Visual companion PDF written: {non_text_pdf_path} "
                      f"({size_bytes:,} bytes)", file=sys.stderr)
                if size_bytes > 52428800:
                    print(f"\nWARNING: Compiled visual bundle PDF ({size_bytes:,} bytes) exceeds "
                          f"the 50 MB size limit (52,428,800 bytes) for Google Cloud multimodal inputs.\n"
                          f"To reduce the size, re-run with lower quality or dimension constraints:\n"
                          f"  omnidoc {args.zip_file} --visual-quality 60 --visual-max-dim 1200\n", file=sys.stderr)
    except Exception as e:
        print(f"error: failed to write output: {e}", file=sys.stderr)
        return 1

    if args.pdf:
        if args.format != FORMAT_MARKDOWN:
            print("warning: --pdf requires markdown format; skipping PDF render.",
                  file=sys.stderr)
        else:
            pdf_path = Path(output_base + ".pdf")
            try:
                render_pdf(content, pdf_path)
                if not args.quiet:
                    print(f"✓ PDF written: {pdf_path}", file=sys.stderr)
            except Exception as e:
                print(f"warning: PDF render failed: {e}", file=sys.stderr)

    return 0


# Console-script entrypoint used by pyproject.toml
def main_cli() -> None:
    sys.exit(main())


if __name__ == "__main__":
    main_cli()
